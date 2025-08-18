import asyncio
from agentpress.tool import ToolResult, openapi_schema, usage_example
from sandbox.tool_base import SandboxToolsBase
from agentpress.thread_manager import ThreadManager
from typing import List, Dict, Optional, Any
import json
import base64
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import tempfile
import os

import hashlib
from PIL import Image
import re
import asyncio
import random
import httpx


class SandboxPresentationToolV2(SandboxToolsBase):
    def __init__(self, project_id: str, thread_manager: ThreadManager):
        super().__init__(project_id, thread_manager)
        self.presentations_dir = "presentations"
        self.images_cache = {}
        
    async def _ensure_presentations_dir(self):
        full_path = f"{self.workspace_path}/{self.presentations_dir}"
        try:
            await self.sandbox.fs.create_folder(full_path, "755")
        except:
            pass
    
    def _get_display_url(self, image_url: str) -> str:
        if not image_url:
            return ""
            
        if image_url.startswith("unsplash:"):
            keyword = image_url.replace("unsplash:", "").strip()
            return f"https://source.unsplash.com/1920x1080/?{keyword}"
        
        return image_url
    
    async def _download_and_cache_image(self, image_url: str, presentation_dir: str) -> Optional[str]:
        """Download an image and cache it locally. Returns the relative path from workspace root."""
        if not image_url:
            return None
            
        # Check if already cached
        if image_url in self.images_cache:
            return self.images_cache[image_url]
        
        try:
            download_url = self._get_display_url(image_url)
            
            # Generate unique filename based on URL
            url_hash = hashlib.md5(image_url.encode()).hexdigest()[:8]
            images_dir = f"{presentation_dir}/images"
            full_images_dir = f"{self.workspace_path}/{images_dir}"
            
            # Ensure images directory exists
            try:
                await self.sandbox.fs.create_folder(full_images_dir, "755")
            except:
                pass
            
            image_data: bytes | None = None
            
            # Try to download the image
            headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
            try:
                async with httpx.AsyncClient(follow_redirects=True, timeout=20.0) as client:
                    resp = await client.get(download_url, headers=headers)
                    resp.raise_for_status()
                    image_data = resp.content
            except Exception as e:
                # Fallback to curl if httpx fails
                try:
                    tmp_path = f"/tmp/img_{url_hash}"
                    cmd = f"/bin/sh -c 'curl -fsSL -A \"Mozilla/5.0\" \"{download_url}\" -o {tmp_path}'"
                    res = await self.sandbox.process.exec(cmd, timeout=30)
                    if getattr(res, "exit_code", 1) == 0:
                        image_data = await self.sandbox.fs.download_file(tmp_path)
                        try:
                            await self.sandbox.process.exec(f"/bin/sh -c 'rm -f {tmp_path}'", timeout=10)
                        except:
                            pass
                except Exception:
                    image_data = None
            
            if not image_data:
                print(f"Failed to download image from {download_url}")
                return None
            
            # Process and save the image
            try:
                img = Image.open(io.BytesIO(image_data))
                output = io.BytesIO()
                
                # Convert to RGB if necessary
                if img.mode in ("RGBA", "LA", "P"):
                    background = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == "P":
                        img = img.convert("RGBA")
                    background.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
                    img = background
                else:
                    img = img.convert("RGB")
                
                img.save(output, format="JPEG", quality=90)
                image_data = output.getvalue()
                filename = f"img_{url_hash}.jpg"
            except Exception:
                # If image processing fails, save as-is
                filename = f"img_{url_hash}.jpg"
            
            # Save the image
            image_path = f"{images_dir}/{filename}"
            full_image_path = f"{self.workspace_path}/{image_path}"
            await self.sandbox.fs.upload_file(image_data, full_image_path)
            
            # Cache and return the relative path
            self.images_cache[image_url] = image_path
            return image_path
            
        except Exception as e:
            print(f"Failed to download image {image_url}: {e}")
            return None

    def _safe_name_variants(self, name: str) -> List[str]:
        """Generate safe name variants for file/folder naming."""
        base = "".join(c if c.isalnum() or c in "-_" else '-' for c in name).lower()
        while "--" in base:
            base = base.replace("--", "-")
        while "__" in base:
            base = base.replace("__", "_")
        hyphen = base.replace("_", "-")
        underscore = base.replace("-", "_")
        variants = []
        for v in [base, hyphen, underscore]:
            if v not in variants:
                variants.append(v)
        return variants
    
    async def _process_slide_images(self, slide: Dict, presentation_dir: str) -> Dict:
        """Process all images in a slide and download them if needed."""
        content = slide.get("content", {})
        
        # Process single image
        if "image" in content:
            image_info = content["image"]
            if isinstance(image_info, str):
                # Convert string to dict format
                local_path = await self._download_and_cache_image(image_info, presentation_dir)
                if local_path:
                    content["image"] = {
                        "url": image_info,
                        "local_path": local_path
                    }
            elif isinstance(image_info, dict) and "url" in image_info:
                # Download if not already downloaded
                if "local_path" not in image_info:
                    local_path = await self._download_and_cache_image(
                        image_info["url"], 
                        presentation_dir
                    )
                    if local_path:
                        image_info["local_path"] = local_path
        
        # Process image grid
        if "images" in content:
            processed_images = []
            for img in content["images"]:
                if isinstance(img, str):
                    local_path = await self._download_and_cache_image(img, presentation_dir)
                    if local_path:
                        processed_images.append({
                            "url": img,
                            "local_path": local_path
                        })
                    else:
                        processed_images.append({"url": img})
                elif isinstance(img, dict):
                    if "url" in img and "local_path" not in img:
                        local_path = await self._download_and_cache_image(
                            img["url"],
                            presentation_dir
                        )
                        if local_path:
                            img["local_path"] = local_path
                    processed_images.append(img)
            content["images"] = processed_images
        
        return slide

    @openapi_schema({
        "type": "function",
        "function": {
            "name": "create_presentation",
            "description": """
            Create a professional presentation using structured JSON format.
            The presentation will be saved as JSON and can be previewed in the browser or exported to PPTX.
            
            IMPORTANT: Generate a structured JSON with specific layout types. Each slide must have:
            - layout: One of 'title', 'title-bullets', 'title-content', 'two-column', 'image-text', 
              'quote', 'section', 'blank', 'hero-image', 'image-grid', 'comparison', 'timeline', 'stats'
            - content: Object with fields specific to that layout type
            
            Images can be specified as:
            - URLs (will be downloaded and embedded)
            - Unsplash search terms using "unsplash:keyword" format
            - Local file paths in the workspace
            
            The tool will handle all formatting and ensure consistency across preview and export.
            
            THEME SELECTION: Instead of a fixed theme, choose a theme that best fits the content:
            - 'corporate-blue': Professional blue theme for business presentations
            - 'modern-purple': Modern purple/pink gradient for tech/startup
            - 'minimal-mono': Black and white minimalist design
            - 'ocean-teal': Calming teal and aqua colors
            - 'sunset-warm': Warm orange and red tones
            - 'forest-green': Natural green palette
            - 'midnight-dark': Dark mode with bright accents
            - 'pastel-soft': Soft pastel colors
            - 'bold-contrast': High contrast bold colors
            - 'elegant-gold': Sophisticated gold and navy
            
            Select the theme that best matches the presentation's content and purpose.
            """,
            "parameters": {
                "type": "object",
                "properties": {
                    "presentation_name": {
                        "type": "string",
                        "description": "Name for the presentation (used for file naming)"
                    },
                    "title": {
                        "type": "string",
                        "description": "Main title of the presentation"
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Optional subtitle or tagline"
                    },
                    "theme": {
                        "type": "string",
                        "enum": ["corporate-blue", "modern-purple", "minimal-mono", "ocean-teal", 
                                "sunset-warm", "forest-green", "midnight-dark", "pastel-soft", 
                                "bold-contrast", "elegant-gold"],
                        "description": "Visual theme - choose based on content and purpose"
                    },
                    "slides": {
                        "type": "array",
                        "description": "Array of slide objects with layout and content",
                        "items": {
                            "type": "object",
                            "properties": {
                                "layout": {
                                    "type": "string",
                                    "enum": ["title", "title-bullets", "title-content", "two-column", "image-text", 
                                            "quote", "section", "blank", "hero-image", "image-grid", 
                                            "comparison", "timeline", "stats"],
                                    "description": "Layout type for the slide"
                                },
                                "content": {
                                    "type": "object",
                                    "description": "Content object specific to the layout type",
                                    "properties": {
                                        "title": {"type": "string"},
                                        "subtitle": {"type": "string"},
                                        "bullets": {"type": "array", "items": {"type": "string"}},
                                        "text": {"type": "string"},
                                        "left_content": {"type": "object"},
                                        "right_content": {"type": "object"},
                                        "image": {
                                            "type": "object",
                                            "properties": {
                                                "url": {"type": "string"},
                                                "alt": {"type": "string"},
                                                "position": {"type": "string", "enum": ["left", "right", "center", "background"]}
                                            }
                                        },
                                        "quote": {"type": "string"},
                                        "author": {"type": "string"},
                                        "notes": {"type": "string"}
                                    }
                                }
                            },
                            "required": ["layout", "content"]
                        }
                    }
                },
                "required": ["presentation_name", "title", "slides"]
            }
        }
    })
    @usage_example('''
        <function_calls>
        <invoke name="create_presentation">
        <parameter name="presentation_name">company_overview</parameter>
        <parameter name="title">Company Overview 2024</parameter>
        <parameter name="subtitle">Innovation Through Technology</parameter>
        <parameter name="theme">modern-purple</parameter>
        <parameter name="slides">[
            {
                "layout": "hero-image",
                "content": {
                    "title": "Welcome to the Future",
                    "subtitle": "Where Innovation Meets Excellence",
                    "image": {
                        "url": "unsplash:technology office",
                        "position": "background"
                    }
                }
            },
            {
                "layout": "stats",
                "content": {
                    "title": "Our Impact in Numbers",
                    "stats": [
                        {"value": "50K+", "label": "Active Users"},
                        {"value": "$2.5M", "label": "Revenue"},
                        {"value": "98%", "label": "Satisfaction"},
                        {"value": "15", "label": "Countries"}
                    ]
                }
            },
            {
                "layout": "image-text",
                "content": {
                    "title": "Our Mission",
                    "text": "We empower businesses with cutting-edge AI technology to transform their operations, enhance productivity, and unlock new possibilities. Our platform combines powerful automation with human creativity.",
                    "image": {
                        "url": "unsplash:artificial intelligence",
                        "position": "right"
                    }
                }
            },
            {
                "layout": "image-grid",
                "content": {
                    "title": "Our Products in Action",
                    "images": [
                        {"url": "unsplash:dashboard analytics", "caption": "Real-time Analytics"},
                        {"url": "unsplash:team collaboration", "caption": "Team Collaboration"},
                        {"url": "unsplash:mobile app", "caption": "Mobile Experience"},
                        {"url": "unsplash:data visualization", "caption": "Data Insights"}
                    ]
                }
            },
            {
                "layout": "title-bullets",
                "content": {
                    "title": "Key Features",
                    "bullets": [
                        "AI-powered automation for repetitive tasks",
                        "Real-time collaboration and communication",
                        "Advanced analytics and reporting",
                        "Enterprise-grade security and compliance",
                        "24/7 customer support and training"
                    ]
                }
            },
            {
                "layout": "quote",
                "content": {
                    "quote": "This platform has transformed how we work. We've saved 40% of our time and increased productivity by 60%.",
                    "author": "Sarah Johnson, CTO at TechCorp"
                }
            },
            {
                "layout": "section",
                "content": {
                    "title": "Let's Build Together",
                    "subtitle": "Start Your Journey Today"
                }
            },
            {
                "layout": "title",
                "content": {
                    "title": "Thank You",
                    "subtitle": "Questions? Let's Connect!",
                    "notes": "Contact us at hello@company.com"
                }
            }
        ]</parameter>
        </invoke>
        </function_calls>
    ''')
    async def create_presentation(
        self,
        presentation_name: str,
        title: str,
        slides: List[Dict[str, Any]],
        subtitle: Optional[str] = None,
        theme: str = "corporate-blue"
    ) -> ToolResult:
        try:
            await self._ensure_sandbox()
            await self._ensure_presentations_dir()
            
            # Validate inputs
            if not presentation_name:
                return self.fail_response("Presentation name is required.")
            
            if not title:
                return self.fail_response("Presentation title is required.")
            
            if not slides or not isinstance(slides, list) or len(slides) == 0:
                return self.fail_response("At least one slide is required.")
            
            # Validate slide structure
            for i, slide in enumerate(slides):
                if 'layout' not in slide:
                    return self.fail_response(f"Slide {i+1} missing 'layout' field")
                if 'content' not in slide:
                    return self.fail_response(f"Slide {i+1} missing 'content' field")
            
            # Create presentation directory
            safe_name = self._safe_name_variants(presentation_name)[0]
            presentation_dir = f"{self.presentations_dir}/{safe_name}"
            full_presentation_path = f"{self.workspace_path}/{presentation_dir}"
            
            try:
                await self.sandbox.fs.create_folder(full_presentation_path, "755")
            except:
                pass
            
            # Process all images in slides (download them during creation)
            download_errors = []
            processed_slides = []
            
            for slide in slides:
                try:
                    processed_slide = await self._process_slide_images(slide.copy(), presentation_dir)
                    processed_slides.append(processed_slide)
                except Exception as e:
                    download_errors.append(f"Error processing slide images: {str(e)}")
                    processed_slides.append(slide)
            
            # Create presentation data
            presentation_data = {
                "version": "2.0",
                "metadata": {
                    "title": title,
                    "subtitle": subtitle,
                    "theme": theme,
                    "created_at": datetime.now().isoformat(),
                    "presentation_name": presentation_name,
                    "total_slides": len(processed_slides)
                },
                "slides": processed_slides,
                "theme_config": self._get_theme_config(theme)
            }
            
            # Save JSON
            json_filename = "presentation.json"
            json_path = f"{presentation_dir}/{json_filename}"
            full_json_path = f"{self.workspace_path}/{json_path}"
            
            json_content = json.dumps(presentation_data, indent=2)
            await self.sandbox.fs.upload_file(
                json_content.encode('utf-8'),
                full_json_path
            )
            
            # Generate HTML preview with downloaded images
            preview_html = self._generate_html_preview(presentation_data)
            preview_path = f"{presentation_dir}/preview.html"
            full_preview_path = f"{self.workspace_path}/{preview_path}"
            
            await self.sandbox.fs.upload_file(
                preview_html.encode('utf-8'),
                full_preview_path
            )
            
            result = {
                "message": f"Successfully created presentation '{title}' with {theme} theme",
                "presentation_name": safe_name,
                "json_file": json_path,
                "preview_url": f"/workspace/{preview_path}",
                "total_slides": len(processed_slides),
                "theme": theme,
                "slides": [
                    {
                        "slide_number": i + 1,
                        "layout": slide["layout"],
                        "title": slide["content"].get("title", f"Slide {i+1}")
                    }
                    for i, slide in enumerate(processed_slides)
                ]
            }
            
            if download_errors:
                result["warnings"] = download_errors
            
            return self.success_response(result)
            
        except Exception as e:
            return self.fail_response(f"Failed to create presentation: {str(e)}")
    
    def _get_theme_config(self, theme: str) -> Dict[str, Any]:
        themes = {
            "corporate-blue": {
                "colors": {
                    "primary": "#1e3a8a",
                    "secondary": "#3b82f6",
                    "accent": "#f59e0b",
                    "background": "#ffffff",
                    "text": "#1e293b",
                    "text_light": "#64748b"
                },
                "fonts": {
                    "heading": "Arial, sans-serif",
                    "body": "Arial, sans-serif"
                }
            },
            "modern-purple": {
                "colors": {
                    "primary": "#7c3aed",
                    "secondary": "#ec4899",
                    "accent": "#06b6d4",
                    "background": "#ffffff",
                    "text": "#1e293b",
                    "text_light": "#64748b"
                },
                "fonts": {
                    "heading": "Inter, sans-serif",
                    "body": "Inter, sans-serif"
                }
            },
            "minimal-mono": {
                "colors": {
                    "primary": "#000000",
                    "secondary": "#525252",
                    "accent": "#000000",
                    "background": "#ffffff",
                    "text": "#000000",
                    "text_light": "#737373"
                },
                "fonts": {
                    "heading": "Helvetica Neue, sans-serif",
                    "body": "Helvetica Neue, sans-serif"
                }
            },
            "ocean-teal": {
                "colors": {
                    "primary": "#0d9488",
                    "secondary": "#06b6d4",
                    "accent": "#0284c7",
                    "background": "#ffffff",
                    "text": "#134e4a",
                    "text_light": "#5eead4"
                },
                "fonts": {
                    "heading": "Roboto, sans-serif",
                    "body": "Roboto, sans-serif"
                }
            },
            "sunset-warm": {
                "colors": {
                    "primary": "#dc2626",
                    "secondary": "#f97316",
                    "accent": "#fbbf24",
                    "background": "#ffffff",
                    "text": "#7c2d12",
                    "text_light": "#ea580c"
                },
                "fonts": {
                    "heading": "Poppins, sans-serif",
                    "body": "Open Sans, sans-serif"
                }
            },
            "forest-green": {
                "colors": {
                    "primary": "#14532d",
                    "secondary": "#16a34a",
                    "accent": "#84cc16",
                    "background": "#ffffff",
                    "text": "#14532d",
                    "text_light": "#22c55e"
                },
                "fonts": {
                    "heading": "Montserrat, sans-serif",
                    "body": "Source Sans Pro, sans-serif"
                }
            },
            "midnight-dark": {
                "colors": {
                    "primary": "#f3f4f6",
                    "secondary": "#60a5fa",
                    "accent": "#f472b6",
                    "background": "#111827",
                    "text": "#f9fafb",
                    "text_light": "#d1d5db"
                },
                "fonts": {
                    "heading": "Space Grotesk, sans-serif",
                    "body": "Inter, sans-serif"
                }
            },
            "pastel-soft": {
                "colors": {
                    "primary": "#c084fc",
                    "secondary": "#fda4af",
                    "accent": "#86efac",
                    "background": "#fef3c7",
                    "text": "#451a03",
                    "text_light": "#92400e"
                },
                "fonts": {
                    "heading": "Quicksand, sans-serif",
                    "body": "Nunito, sans-serif"
                }
            },
            "bold-contrast": {
                "colors": {
                    "primary": "#dc2626",
                    "secondary": "#000000",
                    "accent": "#fbbf24",
                    "background": "#ffffff",
                    "text": "#000000",
                    "text_light": "#525252"
                },
                "fonts": {
                    "heading": "Bebas Neue, sans-serif",
                    "body": "Roboto, sans-serif"
                }
            },
            "elegant-gold": {
                "colors": {
                    "primary": "#1e293b",
                    "secondary": "#f59e0b",
                    "accent": "#92400e",
                    "background": "#fffbeb",
                    "text": "#1e293b",
                    "text_light": "#475569"
                },
                "fonts": {
                    "heading": "Playfair Display, serif",
                    "body": "Lato, sans-serif"
                }
            }
        }
        
        return themes.get(theme, themes["corporate-blue"])
    
    def _generate_html_preview(self, presentation_data: Dict) -> str:
        theme = presentation_data["theme_config"]
        colors = theme["colors"]
        fonts = theme["fonts"]
        
        is_dark = colors["background"].startswith("#1") or colors["background"].startswith("#0")
        
        slides_html = []
        for i, slide in enumerate(presentation_data["slides"]):
            slide_html = self._render_slide_html(slide, i + 1, theme)
            slides_html.append(slide_html)
        
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{presentation_data['metadata']['title']}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: {fonts['body']};
            background: {'#1a1a1a' if is_dark else '#f5f5f5'};
            color: {colors['text']};
            padding: 20px;
        }}
        
        .presentation-header {{
            text-align: center;
            margin-bottom: 30px;
            padding: 20px;
            background: {colors['background']};
            border-radius: 10px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        
        .presentation-title {{
            font-family: {fonts['heading']};
            font-size: 32px;
            color: {colors['primary']};
            margin-bottom: 10px;
        }}
        
        .presentation-subtitle {{
            font-size: 18px;
            color: {colors['text_light']};
        }}
        
        .slides-container {{
            max-width: 1280px;
            margin: 0 auto;
        }}
        
        .slide {{
            width: 100%;
            aspect-ratio: 16/9;
            background: {colors['background']};
            margin-bottom: 30px;
            border-radius: 10px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.1);
            padding: 60px;
            display: flex;
            flex-direction: column;
            position: relative;
            overflow: hidden;
        }}
        
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            color: {colors['text_light']};
            font-size: 14px;
            opacity: 0.6;
        }}
        
        /* Layout-specific styles */
        .layout-title {{
            justify-content: center;
            align-items: center;
            text-align: center;
        }}
        
        .layout-title h1 {{
            font-family: {fonts['heading']};
            font-size: 72px;
            color: {colors['primary']};
            margin-bottom: 30px;
            font-weight: 700;
        }}
        
        .layout-title .subtitle {{
            font-size: 32px;
            color: {colors['secondary']};
            font-weight: 300;
        }}
        
        .layout-title-bullets h2 {{
            font-family: {fonts['heading']};
            font-size: 48px;
            color: {colors['primary']};
            margin-bottom: 40px;
            font-weight: 600;
        }}
        
        .layout-title-bullets ul {{
            list-style: none;
            padding-left: 0;
        }}
        
        .layout-title-bullets li {{
            font-size: 24px;
            margin-bottom: 20px;
            padding-left: 40px;
            position: relative;
            color: {colors['text']};
            line-height: 1.5;
        }}
        
        .layout-title-bullets li:before {{
            content: "●";
            position: absolute;
            left: 0;
            color: {colors['accent']};
            font-size: 24px;
        }}
        
        .layout-two-column {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 60px;
        }}
        
        .layout-two-column > h2 {{
            grid-column: 1 / -1;
            font-family: {fonts['heading']};
            font-size: 48px;
            color: {colors['primary']};
            margin-bottom: 20px;
        }}
        
        .column h3 {{
            font-size: 28px;
            color: {colors['secondary']};
            margin-bottom: 20px;
        }}
        
        .layout-quote {{
            justify-content: center;
            align-items: center;
            text-align: center;
            position: relative;
        }}
        
        .layout-quote.with-overlay {{
            padding: 0;
        }}
        
        .layout-quote .overlay {{
            position: absolute;
            inset: 0;
            background: rgba(0,0,0,0.5);
            border-radius: 10px;
        }}
        
        .layout-quote .quote-content {{
            padding: 60px;
            position: relative;
            z-index: 1;
        }}
        
        .layout-quote .quote-text {{
            font-size: 36px;
            font-style: italic;
            color: {colors['primary']};
            margin-bottom: 30px;
            line-height: 1.5;
            font-weight: 300;
        }}
        
        .layout-quote .quote-author {{
            font-size: 24px;
            color: {colors['secondary']};
        }}
        
        .layout-section {{
            justify-content: center;
            align-items: center;
            text-align: center;
            background: {colors['primary']};
            color: {colors['background']};
        }}
        
        .layout-section h1 {{
            font-family: {fonts['heading']};
            font-size: 72px;
            margin-bottom: 20px;
            font-weight: 700;
        }}
        
        .layout-section .subtitle {{
            font-size: 32px;
            opacity: 0.9;
            font-weight: 300;
        }}
        
        .layout-title-content h2 {{
            font-family: {fonts['heading']};
            font-size: 48px;
            color: {colors['primary']};
            margin-bottom: 40px;
        }}
        
        .layout-title-content .content-text {{
            font-size: 24px;
            line-height: 1.8;
            color: {colors['text']};
        }}
        
        /* Image layouts */
        .layout-image-text {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 60px;
            align-items: center;
        }}
        
        .layout-image-text.image-right {{
            grid-template-columns: 1fr 1fr;
        }}
        
        .layout-image-text.image-left {{
            grid-template-columns: 1fr 1fr;
            direction: rtl;
        }}
        
        .layout-image-text .text-content {{
            direction: ltr;
        }}
        
        .layout-image-text img {{
            width: 100%;
            height: auto;
            border-radius: 10px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.1);
        }}
        
        .layout-hero-image {{
            position: relative;
            padding: 0;
            background-size: cover;
            background-position: center;
            display: flex;
            align-items: center;
            justify-content: center;
        }}
        
        .layout-hero-image .overlay {{
            position: absolute;
            inset: 0;
            background: rgba(0,0,0,0.5);
        }}
        
        .layout-hero-image .content {{
            position: relative;
            z-index: 1;
            text-align: center;
            color: white;
            padding: 60px;
        }}
        
        .layout-hero-image h1 {{
            font-size: 72px;
            margin-bottom: 30px;
            text-shadow: 2px 2px 4px rgba(0,0,0,0.5);
            font-weight: 700;
        }}
        
        .layout-hero-image .subtitle {{
            font-size: 32px;
            opacity: 0.95;
            font-weight: 300;
        }}
        
        .layout-image-grid {{
            padding: 40px;
        }}
        
        .layout-image-grid h2 {{
            font-size: 48px;
            margin-bottom: 40px;
            text-align: center;
            color: {colors['primary']};
        }}
        
        .layout-image-grid .grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 30px;
        }}
        
        .layout-image-grid .grid-item {{
            position: relative;
            border-radius: 10px;
            overflow: hidden;
            box-shadow: 0 5px 20px rgba(0,0,0,0.1);
        }}
        
        .layout-image-grid img {{
            width: 100%;
            height: 250px;
            object-fit: cover;
        }}
        
        .layout-image-grid .caption {{
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 15px;
            font-size: 16px;
        }}
        
        .layout-stats {{
            padding: 60px;
        }}
        
        .layout-stats h2 {{
            font-size: 48px;
            margin-bottom: 60px;
            text-align: center;
            color: {colors['primary']};
        }}
        
        .layout-stats .stats-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
            gap: 40px;
            text-align: center;
        }}
        
        .layout-stats .stat-item {{
            padding: 30px;
            background: {'rgba(255,255,255,0.1)' if is_dark else 'rgba(0,0,0,0.03)'};
            border-radius: 15px;
            border: 2px solid {colors['accent']};
        }}
        
        .layout-stats .stat-value {{
            font-size: 48px;
            font-weight: bold;
            color: {colors['accent']};
            margin-bottom: 10px;
        }}
        
        .layout-stats .stat-label {{
            font-size: 18px;
            color: {colors['text_light']};
        }}
        
        .navigation {{
            position: fixed;
            bottom: 20px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 10px;
            background: {colors['background']};
            padding: 10px 20px;
            border-radius: 50px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.2);
            z-index: 1000;
        }}
        
        .nav-button {{
            padding: 8px 16px;
            background: {colors['primary']};
            color: {colors['background'] if colors['primary'] != '#000000' else '#ffffff'};
            border: none;
            border-radius: 5px;
            cursor: pointer;
            font-size: 14px;
            transition: all 0.3s;
        }}
        
        .nav-button:hover {{
            background: {colors['secondary']};
            transform: translateY(-2px);
        }}
        
        .nav-button:disabled {{
            background: #ccc;
            cursor: not-allowed;
        }}
        
        @media print {{
            .navigation {{
                display: none;
            }}
            
            .slide {{
                page-break-after: always;
                box-shadow: none;
                margin-bottom: 0;
            }}
        }}
    </style>
</head>
<body>
    <div class="presentation-header">
        <h1 class="presentation-title">{presentation_data['metadata']['title']}</h1>
        {f'<p class="presentation-subtitle">{presentation_data["metadata"]["subtitle"]}</p>' if presentation_data['metadata'].get('subtitle') else ''}
    </div>
    
    <div class="slides-container">
        {''.join(slides_html)}
    </div>
    
    <div class="navigation">
        <button class="nav-button" onclick="window.print()">Print/PDF</button>
        <button class="nav-button" onclick="scrollToSlide(0)">First</button>
        <button class="nav-button" onclick="scrollToPrevious()">Previous</button>
        <button class="nav-button" onclick="scrollToNext()">Next</button>
        <button class="nav-button" onclick="scrollToSlide(-1)">Last</button>
    </div>
    
    <script>
        let currentSlide = 0;
        const slides = document.querySelectorAll('.slide');
        
        function scrollToSlide(index) {{
            if (index === -1) index = slides.length - 1;
            if (index >= 0 && index < slides.length) {{
                slides[index].scrollIntoView({{ behavior: 'smooth', block: 'center' }});
                currentSlide = index;
            }}
        }}
        
        function scrollToNext() {{
            if (currentSlide < slides.length - 1) {{
                scrollToSlide(currentSlide + 1);
            }}
        }}
        
        function scrollToPrevious() {{
            if (currentSlide > 0) {{
                scrollToSlide(currentSlide - 1);
            }}
        }}
        
        // Keyboard navigation
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowRight' || e.key === ' ') scrollToNext();
            if (e.key === 'ArrowLeft') scrollToPrevious();
            if (e.key === 'Home') scrollToSlide(0);
            if (e.key === 'End') scrollToSlide(-1);
        }});
    </script>
</body>
</html>"""
        
        return html
    
    def _render_slide_html(self, slide: Dict, slide_number: int, theme: Dict) -> str:
        """Render a slide to HTML using local image paths."""
        layout = slide["layout"]
        content = slide["content"]
        
        # Helper function to get image URL for HTML
        def get_html_image_url(image_info):
            if isinstance(image_info, dict):
                if 'local_path' in image_info and image_info['local_path']:
                    # Use the local downloaded image
                    return f"/workspace/{image_info['local_path']}"
                elif 'url' in image_info:
                    # Fallback to original URL if local path not available
                    return self._get_display_url(image_info['url'])
            elif isinstance(image_info, str):
                return self._get_display_url(image_info)
            return ""
        
        if layout == "title":
            return f"""
        <div class="slide layout-title">
            <h1>{content.get('title', '')}</h1>
            {f'<p class="subtitle">{content["subtitle"]}</p>' if content.get('subtitle') else ''}
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "title-bullets":
            bullets_html = '\n'.join([f'<li>{bullet}</li>' for bullet in content.get('bullets', [])])
            return f"""
        <div class="slide layout-title-bullets">
            <h2>{content.get('title', '')}</h2>
            <ul>{bullets_html}</ul>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "two-column":
            left = content.get('left_content', {})
            right = content.get('right_content', {})
            
            left_bullets = '\n'.join([f'<li>{b}</li>' for b in left.get('bullets', [])])
            right_bullets = '\n'.join([f'<li>{b}</li>' for b in right.get('bullets', [])])
            
            return f"""
        <div class="slide layout-two-column">
            <h2>{content.get('title', '')}</h2>
            <div class="column">
                {f'<h3>{left.get("subtitle", "")}</h3>' if left.get('subtitle') else ''}
                {f'<ul>{left_bullets}</ul>' if left_bullets else ''}
                {f'<p>{left.get("text", "")}</p>' if left.get('text') else ''}
            </div>
            <div class="column">
                {f'<h3>{right.get("subtitle", "")}</h3>' if right.get('subtitle') else ''}
                {f'<ul>{right_bullets}</ul>' if right_bullets else ''}
                {f'<p>{right.get("text", "")}</p>' if right.get('text') else ''}
            </div>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "quote":
            image_info = content.get('image', {})
            image_url = get_html_image_url(image_info)
            
            style = f'background-image: url({image_url}); background-size: cover; background-position: center;' if image_url else ''
            overlay_class = 'with-overlay' if image_url else ''
            
            return f"""
        <div class="slide layout-quote {overlay_class}" style="{style}">
            {f'<div class="overlay"></div>' if image_url else ''}
            <div class="quote-content" style="{'position: relative; z-index: 1;' if image_url else ''}">
                <blockquote class="quote-text" style="{'color: white;' if image_url else ''}">"{content.get('quote', '')}"</blockquote>
                {f'<p class="quote-author" style="{"color: white;" if image_url else ""}">— {content["author"]}</p>' if content.get('author') else ''}
            </div>
            <span class="slide-number" style="{'color: white;' if image_url else ''}">{slide_number}</span>
        </div>"""
        
        elif layout == "section":
            return f"""
        <div class="slide layout-section">
            <h1>{content.get('title', '')}</h1>
            {f'<p class="subtitle">{content["subtitle"]}</p>' if content.get('subtitle') else ''}
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "title-content":
            return f"""
        <div class="slide layout-title-content">
            <h2>{content.get('title', '')}</h2>
            <p class="content-text">{content.get('text', '')}</p>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "image-text":
            image_info = content.get('image', {})
            image_url = get_html_image_url(image_info)
            image_position = 'right'
            
            if isinstance(image_info, dict):
                image_position = image_info.get('position', 'right')
            
            return f"""
        <div class="slide layout-image-text image-{image_position}">
            <div class="text-content">
                <h2>{content.get('title', '')}</h2>
                <p style="font-size: 24px; line-height: 1.6;">{content.get('text', '')}</p>
            </div>
            <div class="image-content">
                {f'<img src="{image_url}" alt="{content.get("title", "")}" />' if image_url else '<div style="background: #f0f0f0; width: 100%; height: 100%; display: flex; align-items: center; justify-content: center; color: #999;">Image not available</div>'}
            </div>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "hero-image":
            image_info = content.get('image', {})
            image_url = get_html_image_url(image_info)
            
            style = f'background-image: url({image_url}); background-size: cover; background-position: center;' if image_url else 'background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);'
            
            return f"""
        <div class="slide layout-hero-image" style="{style}">
            <div class="overlay"></div>
            <div class="content">
                <h1>{content.get('title', '')}</h1>
                {f'<p class="subtitle">{content["subtitle"]}</p>' if content.get('subtitle') else ''}
            </div>
            <span class="slide-number" style="color: white;">{slide_number}</span>
        </div>"""
        
        elif layout == "image-grid":
            images = content.get('images', [])
            grid_html = []
            
            for img in images:
                img_url = get_html_image_url(img)
                caption = img.get('caption', '') if isinstance(img, dict) else ''
                
                if img_url:
                    grid_html.append(f"""
                    <div class="grid-item">
                        <img src="{img_url}" alt="{caption}" />
                        {f'<div class="caption">{caption}</div>' if caption else ''}
                    </div>""")
                else:
                    grid_html.append(f"""
                    <div class="grid-item">
                        <div style="background: #f0f0f0; width: 100%; height: 250px; display: flex; align-items: center; justify-content: center; color: #999;">Image not available</div>
                        {f'<div class="caption">{caption}</div>' if caption else ''}
                    </div>""")
            
            return f"""
        <div class="slide layout-image-grid">
            <h2>{content.get('title', '')}</h2>
            <div class="grid">
                {''.join(grid_html)}
            </div>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        elif layout == "stats":
            stats = content.get('stats', [])
            stats_html = []
            
            for stat in stats:
                if isinstance(stat, dict):
                    stats_html.append(f"""
                    <div class="stat-item">
                        <div class="stat-value">{stat.get('value', '')}</div>
                        <div class="stat-label">{stat.get('label', '')}</div>
                    </div>""")
            
            return f"""
        <div class="slide layout-stats">
            <h2>{content.get('title', '')}</h2>
            <div class="stats-grid">
                {''.join(stats_html)}
            </div>
            <span class="slide-number">{slide_number}</span>
        </div>"""
        
        else:
            return f"""
        <div class="slide layout-blank">
            <span class="slide-number">{slide_number}</span>
        </div>"""

    @openapi_schema({
        "type": "function",
        "function": {
            "name": "export_presentation",
            "description": "Export a JSON presentation to PPTX format with professional formatting.",
            "parameters": {
                "type": "object",
                "properties": {
                    "presentation_name": {
                        "type": "string",
                        "description": "Name of the presentation to export"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["pptx"],
                        "description": "Export format (currently only PPTX supported)"
                    }
                },
                "required": ["presentation_name"]
            }
        }
    })
    async def export_presentation(
        self,
        presentation_name: str,
        format: str = "pptx"
    ) -> ToolResult:
        try:
            await self._ensure_sandbox()
            
            # Resolve the existing presentation directory
            resolved_name = None
            last_error = None
            for candidate in self._safe_name_variants(presentation_name):
                try:
                    json_path_try = f"{self.presentations_dir}/{candidate}/presentation.json"
                    full_json_path_try = f"{self.workspace_path}/{json_path_try}"
                    json_bytes = await self.sandbox.fs.download_file(full_json_path_try)
                    resolved_name = candidate
                    json_content = json_bytes.decode('utf-8')
                    presentation_data = json.loads(json_content)
                    break
                except Exception as e:
                    last_error = e
                    continue
            
            if resolved_name is None:
                return self.fail_response(f"Presentation '{presentation_name}' not found. Tried: {', '.join(self._safe_name_variants(presentation_name))}. Last error: {str(last_error)}")
            
            if format.lower() != "pptx":
                return self.fail_response(f"Format '{format}' not supported. Only 'pptx' is supported.")
            
            # Images should already be downloaded, but check and download any missing ones
            presentation_dir = f"{self.presentations_dir}/{resolved_name}"
            download_errors = []
            
            for slide in presentation_data["slides"]:
                content = slide.get("content", {})
                
                # Check for missing images and download them
                if "image" in content:
                    image_info = content["image"]
                    if isinstance(image_info, dict) and "url" in image_info and "local_path" not in image_info:
                        local_path = await self._download_and_cache_image(
                            image_info["url"], 
                            presentation_dir
                        )
                        if local_path:
                            image_info["local_path"] = local_path
                        else:
                            download_errors.append(f"Failed to download image: {image_info.get('url', 'unknown')}")
                
                if "images" in content:
                    for img in content["images"]:
                        if isinstance(img, dict) and "url" in img and "local_path" not in img:
                            local_path = await self._download_and_cache_image(
                                img["url"],
                                presentation_dir
                            )
                            if local_path:
                                img["local_path"] = local_path
                            else:
                                download_errors.append(f"Failed to download image: {img.get('url', 'unknown')}")
            
            # Create PPTX
            pptx_bytes = await self._create_pptx_from_json(presentation_data)
            
            pptx_filename = f"{resolved_name}.pptx"
            pptx_path = f"{self.presentations_dir}/{resolved_name}/{pptx_filename}"
            full_pptx_path = f"{self.workspace_path}/{pptx_path}"
            
            await self.sandbox.fs.upload_file(pptx_bytes, full_pptx_path)
            
            result = {
                "message": f"Successfully exported presentation to PPTX",
                "export_file": f"/workspace/{pptx_path}",
                "format": "pptx",
                "file_size": len(pptx_bytes),
                "presentation_name": resolved_name
            }
            
            if download_errors:
                result["warnings"] = download_errors
                result["message"] += f" (Note: {len(download_errors)} images could not be downloaded and will be missing from the PPTX)"
            
            return self.success_response(result)
        except Exception as e:
            return self.fail_response(f"Failed to export presentation: {str(e)}")
    
    async def _create_pptx_from_json(self, presentation_data: Dict) -> bytes:
        prs = Presentation()
        
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        theme = presentation_data["theme_config"]
        colors = theme["colors"]
        
        for slide_data in presentation_data["slides"]:
            layout = slide_data["layout"]
            content = slide_data["content"]
            
            if layout == "title":
                self._add_title_slide(prs, content, colors)
            elif layout == "title-bullets":
                self._add_bullets_slide(prs, content, colors)
            elif layout == "two-column":
                self._add_two_column_slide(prs, content, colors)
            elif layout == "quote":
                self._add_quote_slide(prs, content, colors)
            elif layout == "section":
                self._add_section_slide(prs, content, colors)
            elif layout == "title-content":
                self._add_content_slide(prs, content, colors)
            elif layout == "image-text":
                await self._add_image_text_slide_async(prs, content, colors)
            elif layout == "hero-image":
                await self._add_hero_image_slide_async(prs, content, colors)
            elif layout == "image-grid":
                await self._add_image_grid_slide_async(prs, content, colors)
            elif layout == "stats":
                self._add_stats_slide(prs, content, colors)
            else:
                self._add_blank_slide(prs, colors)
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    def _add_title_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title = slide.shapes.title
        title.text = content.get("title", "")
        self._format_text(title.text_frame.paragraphs[0], 72, True, colors["primary"])
        
        if content.get("subtitle") and hasattr(slide.placeholders, '1'):
            subtitle = slide.placeholders[1]
            subtitle.text = content["subtitle"]
            self._format_text(subtitle.text_frame.paragraphs[0], 32, False, colors["secondary"])
    
    def _add_bullets_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title = slide.shapes.title
        title.text = content.get("title", "")
        self._format_text(title.text_frame.paragraphs[0], 48, True, colors["primary"])
        
        if content.get("bullets"):
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            
            for i, bullet in enumerate(content["bullets"]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0
                self._format_text(p, 24, False, colors["text"])
    
    def _add_two_column_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title = slide.shapes.title
        title.text = content.get("title", "")
        self._format_text(title.text_frame.paragraphs[0], 48, True, colors["primary"])
        
        left_content = content.get("left_content", {})
        if len(slide.placeholders) > 1:
            left_box = slide.placeholders[1]
            tf = left_box.text_frame
            tf.clear()
            
            if left_content.get("subtitle"):
                p = tf.paragraphs[0]
                p.text = left_content["subtitle"]
                self._format_text(p, 28, True, colors["secondary"])
            
            if left_content.get("bullets"):
                for bullet in left_content["bullets"]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    self._format_text(p, 20, False, colors["text"])
        
        right_content = content.get("right_content", {})
        if len(slide.placeholders) > 2:
            right_box = slide.placeholders[2]
            tf = right_box.text_frame
            tf.clear()
            
            if right_content.get("subtitle"):
                p = tf.paragraphs[0]
                p.text = right_content["subtitle"]
                self._format_text(p, 28, True, colors["secondary"])
            
            if right_content.get("bullets"):
                for bullet in right_content["bullets"]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    self._format_text(p, 20, False, colors["text"])
    
    def _add_quote_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        left = Inches(2)
        top = Inches(2)
        width = Inches(9.333)
        height = Inches(3)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f'"{content.get("quote", "")}"'
        p.alignment = PP_ALIGN.CENTER
        self._format_text(p, 36, False, colors["primary"], italic=True)
        
        if content.get("author"):
            author_box = slide.shapes.add_textbox(
                Inches(2), Inches(5), Inches(9.333), Inches(1)
            )
            p = author_box.text_frame.paragraphs[0]
            p.text = f"— {content['author']}"
            p.alignment = PP_ALIGN.CENTER
            self._format_text(p, 24, False, colors["secondary"])
    
    def _add_section_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["primary"])
        
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(2)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = content.get("title", "")
        p.alignment = PP_ALIGN.CENTER
        self._format_text(p, 72, True, colors["background"])
        
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(11.333), Inches(1)
            )
            p = subtitle_box.text_frame.paragraphs[0]
            p.text = content["subtitle"]
            p.alignment = PP_ALIGN.CENTER
            self._format_text(p, 32, False, colors["background"])
    
    def _add_content_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title = slide.shapes.title
        title.text = content.get("title", "")
        self._format_text(title.text_frame.paragraphs[0], 48, True, colors["primary"])
        
        if content.get("text"):
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = content["text"]
            p.space_after = Pt(12)
            self._format_text(p, 24, False, colors["text"])
            tf.word_wrap = True
    
    def _add_blank_slide(self, prs, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, colors["background"])
    
    async def _add_image_to_slide(self, slide, image_path: str, left, top, width=None, height=None):
        try:
            if not image_path:
                return None
                
            full_path = f"{self.workspace_path}/{image_path}"
            
            try:
                file_info = await self.sandbox.fs.get_file_info(full_path)
                if file_info.is_dir:
                    print(f"Path is a directory, not an image: {image_path}")
                    return None
            except:
                print(f"Image file not found: {image_path}")
                return None
            
            image_data = await self.sandbox.fs.download_file(full_path)
            
            ext = image_path.split('.')[-1] if '.' in image_path else 'jpg'
            with tempfile.NamedTemporaryFile(suffix=f'.{ext}', delete=False) as tmp_img:
                tmp_img.write(image_data)
                tmp_img.flush()
                
                try:
                    if width and height:
                        pic = slide.shapes.add_picture(tmp_img.name, left, top, width, height)
                    elif width:
                        pic = slide.shapes.add_picture(tmp_img.name, left, top, width=width)
                    elif height:
                        pic = slide.shapes.add_picture(tmp_img.name, left, top, height=height)
                    else:
                        pic = slide.shapes.add_picture(tmp_img.name, left, top)
                    
                    os.unlink(tmp_img.name)
                    return pic
                except Exception as e:
                    os.unlink(tmp_img.name)
                    print(f"Failed to add picture to slide: {e}")
                    return None
                
        except Exception as e:
            print(f"Failed to add image to slide: {e}")
            return None
    
    async def _add_image_text_slide_async(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = content.get("title", "")
        self._format_text(p, 36, True, colors["primary"])
        
        image_info = content.get("image", {})
        image_position = 'right'
        if isinstance(image_info, dict):
            image_position = image_info.get('position', 'right')
        
        if image_position == 'right':
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(5.5), Inches(4.5))
        else:
            text_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(4.5))
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = content.get("text", "")
        self._format_text(p, 20, False, colors["text"])
        
        if isinstance(image_info, dict) and "local_path" in image_info:
            if image_position == 'right':
                await self._add_image_to_slide(
                    slide, 
                    image_info["local_path"],
                    Inches(7), Inches(2),
                    width=Inches(5.5)
                )
            else:
                await self._add_image_to_slide(
                    slide, 
                    image_info["local_path"],
                    Inches(0.5), Inches(2),
                    width=Inches(5.5)
                )
    
    async def _add_hero_image_slide_async(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        image_added = False
        
        image_info = content.get("image", {})
        if isinstance(image_info, dict) and "local_path" in image_info:
            pic = await self._add_image_to_slide(
                slide,
                image_info["local_path"],
                Inches(0), Inches(0),
                width=Inches(13.333), height=Inches(7.5)
            )
            if pic:
                image_added = True
                try:
                    slide.shapes._spTree.remove(pic._element)
                    slide.shapes._spTree.insert(2, pic._element)
                except:
                    pass
        
        if image_added:
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(2), Inches(2),
                Inches(9.333), Inches(3.5)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(0, 0, 0)
            rect.fill.transparency = 0.4
            rect.line.fill.background()
            
            title_color = "#FFFFFF"
            subtitle_color = "#FFFFFF"
        else:
            self._set_slide_background(slide, colors["primary"])
            title_color = colors["background"]
            subtitle_color = colors["background"]
        
        title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(8.333), Inches(1.5))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = content.get("title", "")
        p.alignment = PP_ALIGN.CENTER
        self._format_text(p, 60, True, title_color)
        
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(Inches(2.5), Inches(4), Inches(8.333), Inches(1))
            p = subtitle_box.text_frame.paragraphs[0]
            p.text = content["subtitle"]
            p.alignment = PP_ALIGN.CENTER
            self._format_text(p, 28, False, subtitle_color)
    
    async def _add_image_grid_slide_async(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = content.get("title", "")
        p.alignment = PP_ALIGN.CENTER
        self._format_text(p, 36, True, colors["primary"])
        
        images = content.get("images", [])
        if images:
            num_images = min(len(images), 6)
            
            if num_images <= 2:
                cols = 2
            elif num_images <= 4:
                cols = 2
            else:
                cols = 3
                
            rows = (num_images + cols - 1) // cols
            
            img_width = Inches(3.5)
            img_height = Inches(2.5)
            h_spacing = Inches(0.5)
            v_spacing = Inches(0.5)
            
            total_width = cols * img_width.inches + (cols - 1) * h_spacing.inches
            total_height = rows * img_height.inches + (rows - 1) * v_spacing.inches
            start_left = (13.333 - total_width) / 2
            start_top = 2 + (5.5 - total_height) / 2
            
            for i, img in enumerate(images[:num_images]):
                if isinstance(img, dict) and "local_path" in img:
                    row = i // cols
                    col = i % cols
                    left = Inches(start_left) + col * (img_width + h_spacing)
                    top = Inches(start_top) + row * (img_height + v_spacing)
                    
                    pic = await self._add_image_to_slide(
                        slide,
                        img["local_path"],
                        left, top,
                        width=img_width, height=img_height
                    )
                    
                    if pic and img.get("caption"):
                        caption_box = slide.shapes.add_textbox(
                            left, top + img_height - Inches(0.5),
                            img_width, Inches(0.5)
                        )
                        p = caption_box.text_frame.paragraphs[0]
                        p.text = img["caption"]
                        p.alignment = PP_ALIGN.CENTER
                        self._format_text(p, 12, False, colors["text"])
                        caption_box.fill.solid()
                        caption_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        caption_box.fill.transparency = 0.2
    
    def _add_stats_slide(self, prs, content: Dict, colors: Dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        self._set_slide_background(slide, colors["background"])
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1))
        p = title_box.text_frame.paragraphs[0]
        p.text = content.get("title", "")
        p.alignment = PP_ALIGN.CENTER
        self._format_text(p, 42, True, colors["primary"])
        
        stats = content.get("stats", [])
        if stats:
            num_stats = min(len(stats), 4)
            stat_width = Inches(2.5)
            stat_height = Inches(2)
            h_spacing = Inches(0.5)
            
            total_width = num_stats * stat_width.inches + (num_stats - 1) * h_spacing.inches
            start_left = (13.333 - total_width) / 2
            top = Inches(3)
            
            for i, stat in enumerate(stats[:num_stats]):
                if isinstance(stat, dict):
                    left = Inches(start_left) + i * (stat_width + h_spacing)
                    
                    stat_box = slide.shapes.add_textbox(left, top, stat_width, stat_height)
                    tf = stat_box.text_frame
                    tf.clear()
                    
                    p = tf.paragraphs[0]
                    p.text = str(stat.get("value", ""))
                    p.alignment = PP_ALIGN.CENTER
                    self._format_text(p, 48, True, colors["accent"])
                    
                    p = tf.add_paragraph()
                    p.text = stat.get("label", "")
                    p.alignment = PP_ALIGN.CENTER
                    self._format_text(p, 18, False, colors["text_light"])
    
    def _set_slide_background(self, slide, color: str):
        if color and color != "transparent":
            try:
                rgb = self._hex_to_rgb(color)
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            except Exception as e:
                print(f"Error setting background: {e}")
    
    def _format_text(self, paragraph, size: int, bold: bool, color: str, italic: bool = False):
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        if color:
            rgb = self._hex_to_rgb(color)
            paragraph.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    
    def _hex_to_rgb(self, hex_color: str) -> tuple:
        if not hex_color or not hex_color.startswith('#'):
            return (0, 0, 0)
        try:
            hex_color = hex_color[1:]
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (0, 0, 0) 