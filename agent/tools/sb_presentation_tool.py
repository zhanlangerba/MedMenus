from agentpress.tool import ToolResult, openapi_schema, usage_example
from sandbox.tool_base import SandboxToolsBase
from agentpress.thread_manager import ThreadManager
from typing import List, Dict, Optional, Union
import json
import os
import base64
from datetime import datetime
import requests
import tempfile
import re
from html import unescape
import io

# New imports for python-pptx and HTML parsing
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from bs4 import BeautifulSoup
import cssutils
import logging

# Suppress cssutils warnings
cssutils.log.setLevel(logging.ERROR)

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    print("PIL/Pillow not available - WEBP images will be skipped in PPTX export")

class SandboxPresentationTool(SandboxToolsBase):
    def __init__(self, project_id: str, thread_manager: ThreadManager):
        super().__init__(project_id, thread_manager)
        self.workspace_path = "/workspace"
        self.presentations_dir = "presentations"

    async def _ensure_presentations_dir(self):
        full_path = f"{self.workspace_path}/{self.presentations_dir}"
        try:
            await self.sandbox.fs.create_folder(full_path, "755")
        except:
            pass

    def _generate_slide_html(self, slide: Dict, slide_number: int, total_slides: int, presentation_title: str, custom_css: Optional[str] = None) -> str:
        """Generate HTML for a single slide - ALWAYS maintains 1920x1080 dimensions"""
        
        if custom_css:
            # Ensure custom CSS includes proper slide dimensions
            if ".slide" in custom_css and "1920px" not in custom_css:
                # Prepend dimension enforcement to custom CSS
                css = """
                /* ENFORCED: Presentation slide dimensions */
                .slide {
                    width: 1920px !important;
                    height: 1080px !important;
                    max-width: 100vw;
                    max-height: 100vh;
                    aspect-ratio: 16/9;
                    transform-origin: center center;
                }
                @media screen and (max-width: 1920px), screen and (max-height: 1080px) {
                    .slide {
                        transform: scale(min(100vw / 1920, 100vh / 1080));
                    }
                }
                """ + custom_css
            else:
                css = custom_css
        else:
            css = """
            * {
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }
            
            body {
                font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'Helvetica Neue', Arial, sans-serif;
                width: 100vw;
                height: 100vh;
                display: flex;
                align-items: center;
                justify-content: center;
                overflow: hidden;
                background: #ffffff;
                color: #000000;
            }
            
            .slide {
                /* CRITICAL: Fixed presentation dimensions 1920x1080 (16:9) */
                width: 1920px !important;
                height: 1080px !important;
                max-width: 100vw;
                max-height: 100vh;
                aspect-ratio: 16/9;
                display: flex;
                flex-direction: column;
                justify-content: center;
                padding: 80px;
                position: relative;
                background: #ffffff;
                /* Scale to fit viewport if needed */
                transform-origin: center center;
            }
            
            /* Auto-scale slide to fit viewport while maintaining aspect ratio */
            @media screen and (max-width: 1920px), screen and (max-height: 1080px) {
                .slide {
                    transform: scale(min(100vw / 1920, 100vh / 1080));
                }
            }
            
            h1 {
                font-size: 72px;
                font-weight: 700;
                line-height: 1.1;
                margin-bottom: 40px;
                color: #000000;
            }
            
            h2 {
                font-size: 48px;
                font-weight: 600;
                line-height: 1.2;
                margin-bottom: 30px;
                color: #333333;
            }
            
            p {
                font-size: 24px;
                line-height: 1.6;
                margin-bottom: 20px;
                color: #333333;
            }
            
            ul {
                list-style: none;
                margin: 30px 0;
                padding-left: 0;
            }
            
            li {
                font-size: 24px;
                line-height: 1.8;
                margin: 15px 0;
                padding-left: 30px;
                position: relative;
                color: #333333;
            }
            
            li::before {
                content: "•";
                position: absolute;
                left: 0;
                color: #000000;
            }
            
            .slide-number {
                position: absolute;
                bottom: 40px;
                right: 40px;
                font-size: 18px;
                color: #666666;
            }
            
            /* Flat design - no shadows, gradients, or animations */
            img {
                max-width: 100%;
                height: auto;
                border: 2px solid #e0e0e0;
            }
            
            .content-section {
                max-width: 100%;
            }
            """
        
        slide_html = slide.get('html', '')
        
        if not slide_html:
            title = slide.get('title', '')
            content = slide.get('content', '')
            
            slide_html = f"""
            <div class="slide">
                <div class="content-section">
                    {f'<h1>{title}</h1>' if title else ''}
                    {content if isinstance(content, str) else ''}
                </div>
                <div class="slide-number">{slide_number} / {total_slides}</div>
            </div>
            """
        
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{presentation_title} - Slide {slide_number}</title>
    <style>
    {css}
    </style>
</head>
<body>
    {slide_html}
</body>
</html>"""
        
        return html

    def _generate_presentation_index(self, title: str, slides: List[str]) -> str:
        slide_links = '\n'.join([
            f'<li><a href="{slide}" target="slide-frame">Slide {i+1}</a></li>'
            for i, slide in enumerate(slides)
        ])
        
        html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title} - Presentation</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            display: flex;
            height: 100vh;
            background: #f5f5f5;
        }}
        .sidebar {{
            width: 250px;
            background: #ffffff;
            border-right: 1px solid #e0e0e0;
            padding: 20px;
            overflow-y: auto;
        }}
        .sidebar h2 {{
            font-size: 20px;
            margin-bottom: 20px;
            color: #333333;
        }}
        .sidebar ul {{
            list-style: none;
        }}
        .sidebar li {{
            margin: 10px 0;
        }}
        .sidebar a {{
            color: #0066cc;
            text-decoration: none;
            display: block;
            padding: 8px 12px;
            border: 1px solid transparent;
            transition: all 0.2s;
        }}
        .sidebar a:hover {{
            background: #f0f0f0;
            border-color: #e0e0e0;
        }}
        .content {{
            flex: 1;
            display: flex;
            align-items: center;
            justify-content: center;
            background: #ffffff;
        }}
        iframe {{
            width: 95%;
            height: 95%;
            border: 1px solid #e0e0e0;
            background: white;
        }}
        .fullscreen-btn {{
            position: fixed;
            top: 20px;
            right: 20px;
            padding: 10px 20px;
            background: #0066cc;
            color: white;
            border: none;
            cursor: pointer;
            font-size: 14px;
            z-index: 1000;
        }}
        .fullscreen-btn:hover {{
            background: #0052a3;
        }}
    </style>
</head>
<body>
    <div class="sidebar">
        <h2>{title}</h2>
        <ul>
            {slide_links}
        </ul>
    </div>
    <div class="content">
        <iframe name="slide-frame" src="{slides[0] if slides else ''}" frameborder="0"></iframe>
    </div>
    <button class="fullscreen-btn" onclick="document.querySelector('iframe').requestFullscreen()">
        Fullscreen
    </button>
</body>
</html>"""
        return html

    @openapi_schema({
        "type": "function",
        "function": {
            "name": "create_presentation",
            "description": "Create a professional presentation by generating raw HTML and CSS for each slide. CRITICAL: Every slide MUST be exactly 1920x1080 pixels (16:9 aspect ratio) - these are standard PowerPoint/presentation dimensions. The agent should create FLAT DESIGN slides with NO gradients, NO shadows, NO animations - just clean, simple, flat colors and typography. Each slide should be self-contained HTML with embedded CSS styling. The .slide class MUST have width: 1920px and height: 1080px.",
            "parameters": {
                "type": "object",
                "properties": {
                    "presentation_name": {
                        "type": "string",
                        "description": "Name of the presentation (used for file naming)"
                    },
                    "title": {
                        "type": "string",
                        "description": "The main title of the presentation"
                    },
                    "slides": {
                        "type": "array",
                        "description": "Array of slides with raw HTML and CSS",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {
                                    "type": "string",
                                    "description": "The title of the slide (for reference)"
                                },
                                "html": {
                                    "type": "string",
                                    "description": "Complete HTML content for the slide. MUST contain a div with class='slide' that will be styled to exactly 1920x1080 pixels. All content should be inside this div. The slide div is your canvas with fixed dimensions of 1920x1080 (16:9 aspect ratio)."
                                },
                                "css": {
                                    "type": "string",
                                    "description": "Custom CSS for this slide. CRITICAL REQUIREMENTS: 1) The .slide class MUST have width: 1920px and height: 1080px. 2) FLAT DESIGN only: Use solid colors, NO gradients, NO box-shadows, NO text-shadows, NO animations. 3) Keep typography clean and spacing consistent. The slide dimensions (1920x1080) are mandatory for proper presentation format."
                                }
                            },
                            "required": ["title", "html", "css"]
                        }
                    }
                },
                "required": ["presentation_name", "title", "slides"]
            }
        }
    })
    @usage_example('''
        <function_calls>
        <invoke name="create_presentation">
        <parameter name="presentation_name">company_overview</parameter>
        <parameter name="title">Company Overview 2024</parameter>
        <parameter name="slides">[
            {
                "title": "Title Slide",
                "html": "<div class='slide title-slide'><h1>Company Overview</h1><p class='subtitle'>Building the Future Together</p><p class='date'>2024</p></div>",
                "css": "* {margin: 0; padding: 0; box-sizing: border-box;} body {font-family: 'Helvetica Neue', Arial, sans-serif; background: #1a1a1a; color: #ffffff; display: flex; align-items: center; justify-content: center; height: 100vh;} .slide {width: 1920px !important; height: 1080px !important; display: flex; flex-direction: column; justify-content: center; align-items: center; text-align: center; padding: 80px;} h1 {font-size: 96px; font-weight: 300; margin-bottom: 30px; letter-spacing: -2px;} .subtitle {font-size: 36px; color: #cccccc; margin-bottom: 60px;} .date {font-size: 24px; color: #999999;}"
            },
            {
                "title": "Our Mission",
                "html": "<div class='slide'><h2>Our Mission</h2><div class='content'><p class='statement'>To empower businesses with innovative solutions that drive growth and success</p><ul><li>Customer-focused approach</li><li>Continuous innovation</li><li>Sustainable practices</li></ul></div></div>",
                "css": "* {margin: 0; padding: 0; box-sizing: border-box;} body {font-family: 'Helvetica Neue', Arial, sans-serif; background: #ffffff; color: #333333; display: flex; align-items: center; justify-content: center; height: 100vh;} .slide {width: 1920px !important; height: 1080px !important; padding: 120px;} h2 {font-size: 72px; font-weight: 600; margin-bottom: 80px; color: #000000;} .statement {font-size: 36px; line-height: 1.5; margin-bottom: 60px; color: #555555;} ul {list-style: none; padding: 0;} li {font-size: 28px; margin: 20px 0; padding-left: 40px; position: relative;} li:before {content: '→'; position: absolute; left: 0; color: #0066cc;}"
            }
        ]</parameter>
        </invoke>
        </function_calls>
    ''')
    async def create_presentation(
        self,
        presentation_name: str,
        title: str,
        slides: List[Dict]
    ) -> ToolResult:
        try:
            await self._ensure_sandbox()
            await self._ensure_presentations_dir()
            
            if not presentation_name:
                return self.fail_response("Presentation name is required.")
            
            if not title:
                return self.fail_response("Presentation title is required.")
            
            if not slides or not isinstance(slides, list):
                return self.fail_response("At least one slide is required.")
            
            safe_name = "".join(c for c in presentation_name if c.isalnum() or c in "-_").lower()
            presentation_dir = f"{self.presentations_dir}/{safe_name}"
            full_presentation_path = f"{self.workspace_path}/{presentation_dir}"
            
            try:
                await self.sandbox.fs.create_folder(full_presentation_path, "755")
            except:
                pass
            
            slide_files = []
            slide_info = []
            
            for i, slide in enumerate(slides, 1):
                custom_css = slide.get('css', '')
                slide_html = self._generate_slide_html(slide, i, len(slides), title, custom_css)
                
                slide_filename = f"slide_{i:02d}.html"
                slide_path = f"{presentation_dir}/{slide_filename}"
                full_slide_path = f"{self.workspace_path}/{slide_path}"
                
                await self.sandbox.fs.upload_file(slide_html.encode(), full_slide_path)
                slide_files.append(slide_filename)
                
                slide_info.append({
                    "slide_number": i,
                    "title": slide.get("title", f"Slide {i}"),
                    "file": slide_path,
                    "preview_url": f"/workspace/{slide_path}"
                })
            
            index_html = self._generate_presentation_index(title, slide_files)
            index_path = f"{presentation_dir}/index.html"
            full_index_path = f"{self.workspace_path}/{index_path}"
            await self.sandbox.fs.upload_file(index_html.encode(), full_index_path)
            
            # Save metadata
            metadata = {
                "presentation_name": presentation_name,
                "title": title,
                "total_slides": len(slides),
                "created_at": datetime.now().isoformat(),
                "slides": slide_info,
                "index_file": index_path,
                "original_slides_data": slides
            }
            
            metadata_path = f"{presentation_dir}/metadata.json"
            full_metadata_path = f"{self.workspace_path}/{metadata_path}"
            await self.sandbox.fs.upload_file(json.dumps(metadata, indent=2).encode(), full_metadata_path)
            
            return self.success_response({
                "message": f"Presentation '{title}' created successfully with {len(slides)} slides using custom HTML/CSS",
                "presentation_path": presentation_dir,
                "index_file": index_path,
                "slides": slide_info,
                "presentation_name": presentation_name,
                "title": title,
                "total_slides": len(slides),
                "note": "Slides created with flat design principles - no gradients, shadows, or animations"
            })
            
        except Exception as e:
            return self.fail_response(f"Failed to create presentation: {str(e)}")

    @openapi_schema({
        "type": "function",
        "function": {
            "name": "export_presentation",
            "description": "Export a presentation to PDF or PPTX format. Note: This requires additional tools to be installed in the environment.",
            "parameters": {
                "type": "object",
                "properties": {
                    "presentation_name": {
                        "type": "string",
                        "description": "Name of the presentation to export"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["pdf", "pptx"],
                        "description": "Export format"
                    }
                },
                "required": ["presentation_name", "format"]
            }
        }
    })
    def _clean_html_text(self, html_text: str) -> str:
        clean = re.compile('<.*?>')
        text = re.sub(clean, '', html_text)
        text = unescape(text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def _hex_to_rgb(self, hex_color: str) -> tuple:
        if not hex_color.startswith('#'):
            return (45, 45, 47)
        try:
            hex_color = hex_color[1:]
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        except:
            return (45, 45, 47)

    async def _download_image_for_pptx(self, url: str) -> Optional[bytes]:
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
            }
            response = requests.get(url, timeout=10, headers=headers)
            response.raise_for_status()
            
            content_type = response.headers.get('Content-Type', '')
            if not content_type.startswith('image/'):
                return None
            
            image_data = response.content
            
            if PIL_AVAILABLE:
                try:
                    with Image.open(io.BytesIO(image_data)) as img:
                        if img.format in ['WEBP'] or img.format not in ['JPEG', 'PNG', 'GIF', 'BMP', 'TIFF']:
                            print(f"Converting image from {img.format} to JPEG for PPTX compatibility")
                            if img.mode in ['RGBA', 'LA', 'P']:
                                background = Image.new('RGB', img.size, (255, 255, 255))
                                if img.mode == 'P':
                                    img = img.convert('RGBA')
                                background.paste(img, mask=img.split()[-1] if img.mode in ['RGBA', 'LA'] else None)
                                img = background
                            elif img.mode != 'RGB':
                                img = img.convert('RGB')
                            
                            output = io.BytesIO()
                            img.save(output, format='JPEG', quality=85)
                            return output.getvalue()
                    
                    return image_data
                    
                except Exception as convert_error:
                    print(f"Error converting image format: {convert_error}")
                    return image_data
            else:
                if 'webp' in content_type.lower():
                    print("WEBP image detected but PIL not available - skipping image")
                    return None
                return image_data
                
        except Exception as e:
            print(f"Error downloading image: {e}")
            return None

    async def _create_pptx_presentation(self, metadata: Dict, slides_data: List[Dict], color_scheme = None) -> bytes:
        """Create a PPTX presentation from HTML slides using python-pptx"""
        prs = Presentation()
        
        # Set slide size to 16:9 (default is already 16:9 in python-pptx)
        prs.slide_width = Inches(13.333)  # 1920 pixels at 144 DPI
        prs.slide_height = Inches(7.5)    # 1080 pixels at 144 DPI
        
        # Process each slide from the metadata
        for i, slide_data in enumerate(slides_data):
            # Get HTML and CSS content from original slides data
            html_content = slide_data.get('html', '')
            css_content = slide_data.get('css', '')
            title = slide_data.get('title', f'Slide {i+1}')
            
            # If we have HTML content, parse it and convert to PPTX
            if html_content or css_content:
                await self._add_html_slide_to_pptx(prs, html_content, css_content, title)
            else:
                # Fallback for old format slides
                await self._add_legacy_slide_to_pptx(prs, slide_data)
        
        # Save presentation to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()

    async def _add_html_slide_to_pptx(self, prs, html_content: str, css_content: str, slide_title: str):
        """Convert HTML/CSS slide to PPTX format using BeautifulSoup and python-pptx"""
        
        # Parse CSS to extract styles
        styles = self._parse_css_styles(css_content) if css_content else {}
        
        # Parse HTML content
        soup = BeautifulSoup(html_content, 'html.parser') if html_content else None
        
        # Add a blank slide with blank layout
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Extract background color from styles
        slide_styles = styles.get('.slide', {})
        bg_color = slide_styles.get('background', '#ffffff')
        if bg_color and bg_color != 'transparent':
            self._set_slide_background(slide, bg_color)
        
        # Find the main slide div
        slide_div = soup.find('div', class_='slide') if soup else None
        
        if slide_div:
            # Process slide content
            await self._process_slide_content(slide, slide_div, styles)
        else:
            # Fallback: just add the title
            self._add_title_to_slide(slide, slide_title)
    
    def _parse_css_styles(self, css_content: str) -> Dict:
        """Parse CSS content and extract styles for each selector"""
        styles = {}
        try:
            sheet = cssutils.parseString(css_content)
            for rule in sheet:
                if hasattr(rule, 'selectorText') and hasattr(rule, 'style'):
                    selector = rule.selectorText
                    rule_styles = {}
                    for prop in rule.style:
                        rule_styles[prop.name] = prop.value
                    styles[selector] = rule_styles
        except Exception as e:
            print(f"Error parsing CSS: {e}")
        return styles
    
    def _set_slide_background(self, slide, color: str):
        """Set slide background color"""
        try:
            rgb = self._hex_to_rgb(color)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        except Exception as e:
            print(f"Error setting background: {e}")
    
    async def _process_slide_content(self, slide, slide_div, styles: Dict):
        """Process HTML content and add to PPTX slide"""
        
        # Track vertical position for elements
        top_position = Inches(0.5)
        
        # Process each element in the slide
        for element in slide_div.children:
            if not hasattr(element, 'name'):
                continue
                
            if element.name == 'h1':
                top_position = self._add_heading(slide, element, styles, 1, top_position)
            elif element.name == 'h2':
                top_position = self._add_heading(slide, element, styles, 2, top_position)
            elif element.name == 'h3':
                top_position = self._add_heading(slide, element, styles, 3, top_position)
            elif element.name == 'p':
                top_position = self._add_paragraph(slide, element, styles, top_position)
            elif element.name == 'ul':
                top_position = self._add_bullet_list(slide, element, styles, top_position)
            elif element.name == 'ol':
                top_position = self._add_numbered_list(slide, element, styles, top_position)
            elif element.name == 'div':
                # Recursively process div contents
                for child in element.children:
                    if hasattr(child, 'name'):
                        if child.name == 'h1':
                            top_position = self._add_heading(slide, child, styles, 1, top_position)
                        elif child.name == 'h2':
                            top_position = self._add_heading(slide, child, styles, 2, top_position)
                        elif child.name == 'p':
                            top_position = self._add_paragraph(slide, child, styles, top_position)
                        elif child.name == 'ul':
                            top_position = self._add_bullet_list(slide, child, styles, top_position)
            elif element.name == 'img':
                top_position = await self._add_image(slide, element, top_position)
    
    def _add_heading(self, slide, element, styles: Dict, level: int, top_position):
        """Add a heading to the slide"""
        text = element.get_text(strip=True)
        if not text:
            return top_position
            
        # Determine font size based on heading level
        font_sizes = {1: 48, 2: 36, 3: 28}
        font_size = font_sizes.get(level, 24)
        
        # Get styles for this heading level
        selector = f'h{level}'
        element_styles = styles.get(selector, {})
        
        # Extract styling
        color = element_styles.get('color', '#000000')
        
        # Add text box
        left = Inches(1)
        width = Inches(11.333)  # Leave margins
        height = Inches(1)
        
        text_box = slide.shapes.add_textbox(left, top_position, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        
        # Set color
        rgb = self._hex_to_rgb(color)
        p.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        # Center align for h1
        if level == 1:
            p.alignment = PP_ALIGN.CENTER
        
        return top_position + height + Inches(0.2)
    
    def _add_paragraph(self, slide, element, styles: Dict, top_position):
        """Add a paragraph to the slide"""
        text = element.get_text(strip=True)
        if not text:
            return top_position
            
        # Get paragraph styles
        element_styles = styles.get('p', {})
        color = element_styles.get('color', '#333333')
        
        # Check for special classes
        if 'class' in element.attrs:
            for cls in element['class']:
                class_styles = styles.get(f'.{cls}', {})
                if 'color' in class_styles:
                    color = class_styles['color']
        
        # Add text box
        left = Inches(1)
        width = Inches(11.333)
        height = Inches(0.8)
        
        text_box = slide.shapes.add_textbox(left, top_position, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        
        rgb = self._hex_to_rgb(color)
        p.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        # Check for center alignment
        if 'subtitle' in element.get('class', []) or 'date' in element.get('class', []):
            p.alignment = PP_ALIGN.CENTER
        
        return top_position + height + Inches(0.1)
    
    def _add_bullet_list(self, slide, element, styles: Dict, top_position):
        """Add a bullet list to the slide"""
        items = element.find_all('li')
        if not items:
            return top_position
            
        # Get list styles
        element_styles = styles.get('li', {})
        color = element_styles.get('color', '#333333')
        
        # Calculate height needed
        height = Inches(0.5 * len(items))
        
        # Add text box
        left = Inches(1.5)
        width = Inches(10.833)
        
        text_box = slide.shapes.add_textbox(left, top_position, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        for i, item in enumerate(items):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = item.get_text(strip=True)
            p.font.size = Pt(16)
            p.level = 0
            
            rgb = self._hex_to_rgb(color)
            p.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        return top_position + height + Inches(0.2)
    
    def _add_numbered_list(self, slide, element, styles: Dict, top_position):
        """Add a numbered list to the slide"""
        items = element.find_all('li')
        if not items:
            return top_position
            
        # Similar to bullet list but with numbers
        height = Inches(0.5 * len(items))
        
        left = Inches(1.5)
        width = Inches(10.833)
        
        text_box = slide.shapes.add_textbox(left, top_position, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        for i, item in enumerate(items, 1):
            p = text_frame.add_paragraph() if i > 1 else text_frame.paragraphs[0]
            p.text = f"{i}. {item.get_text(strip=True)}"
            p.font.size = Pt(16)
            
        return top_position + height + Inches(0.2)
    
    async def _add_image(self, slide, element, top_position):
        """Add an image to the slide"""
        src = element.get('src', '')
        if not src:
            return top_position
            
        try:
            # Download image
            image_data = await self._download_image_for_pptx(src)
            if image_data:
                # Save to temp file
                with tempfile.NamedTemporaryFile(suffix='.jpg', delete=False) as tmp_img:
                    tmp_img.write(image_data)
                    tmp_img.flush()
                    
                    # Add picture to slide
                    left = Inches(2)
                    height = Inches(3)
                    pic = slide.shapes.add_picture(tmp_img.name, left, top_position, height=height)
                    
                    os.unlink(tmp_img.name)
                    return top_position + height + Inches(0.2)
        except Exception as e:
            print(f"Error adding image: {e}")
            
        return top_position
    
    def _add_title_to_slide(self, slide, title: str):
        """Add a simple title to slide as fallback"""
        left = Inches(1)
        top = Inches(3)
        width = Inches(11.333)
        height = Inches(1.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    
    async def _add_legacy_slide_to_pptx(self, prs, slide_data: Dict):
        """Handle legacy slide format (fallback)"""
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        title = slide_data.get('title', '')
        content = slide_data.get('content', {})
        
        # Set background
        bg_color = slide_data.get('background_color', '#ffffff')
        self._set_slide_background(slide, bg_color)
        
        # Add title
        if title:
            self._add_title_to_slide(slide, title)

    async def export_presentation(
        self,
        presentation_name: str,
        format: str = "pptx"
    ) -> ToolResult:
        try:
            await self._ensure_sandbox()
            safe_name = "".join(c for c in presentation_name if c.isalnum() or c in "-_").lower()
            presentation_dir = f"{self.presentations_dir}/{safe_name}"
            
            metadata_path = f"{self.workspace_path}/{presentation_dir}/metadata.json"
            try:
                metadata_content = await self.sandbox.fs.download_file(metadata_path)
                metadata = json.loads(metadata_content.decode())
            except Exception as e:
                return self.fail_response(f"Presentation '{presentation_name}' (safe_name: '{safe_name}') not found at path '{metadata_path}'. Error: {str(e)}")
            
            if format.lower() == "pptx":
                slides_data = metadata.get('original_slides_data', [])
                
                default_bg_color = '#ffffff'
                default_text_color = '#000000'
                
                if not slides_data:
                    slides_data = []
                    for slide_info in metadata.get('slides', []):
                        slides_data.append({
                            'title': slide_info.get('title', f"Slide {slide_info.get('slide_number', 1)}"),
                            'content': {'subtitle': 'Content from HTML slide'},
                            'layout': 'default',
                            'background_color': default_bg_color,
                            'text_color': default_text_color
                        })
                else:
                    for slide in slides_data:
                        if 'background_color' not in slide:
                            slide['background_color'] = default_bg_color
                        if 'text_color' not in slide:
                            slide['text_color'] = default_text_color
                
                try:
                    pptx_data = await self._create_pptx_presentation(metadata, slides_data, None)
                except Exception as e:
                    return self.fail_response(f"PPTX generation failed: {str(e)}")
                
                pptx_filename = f"{safe_name}.pptx"
                pptx_path = f"{presentation_dir}/{pptx_filename}"
                full_pptx_path = f"{self.workspace_path}/{pptx_path}"
                
                print(f"PPTX Debug - safe_name: {safe_name}")
                print(f"PPTX Debug - pptx_filename: {pptx_filename}")
                print(f"PPTX Debug - pptx_path: {pptx_path}")
                print(f"PPTX Debug - full_pptx_path: {full_pptx_path}")
                print(f"PPTX Debug - pptx_data size: {len(pptx_data)} bytes")
                
                await self.sandbox.fs.upload_file(pptx_data, full_pptx_path)
                
                try:
                    file_info = await self.sandbox.fs.get_file_info(full_pptx_path)
                    print(f"PPTX Debug - File created successfully: {file_info.size} bytes")
                except Exception as e:
                    print(f"PPTX Debug - Error verifying file: {str(e)}")
                    return self.fail_response(f"Failed to verify PPTX file creation: {str(e)}")
                
                return self.success_response({
                    "message": f"Presentation exported successfully as PPTX",
                    "export_file": pptx_path,
                    "download_url": f"/workspace/{pptx_path}",
                    "format": "pptx",
                    "presentation_name": presentation_name,
                    "file_size": len(pptx_data)
                })
            else:
                return self.fail_response(f"Export format '{format}' not yet implemented. Only PPTX is currently supported.")
            
        except ImportError:
            return self.fail_response("PPTX export requires 'aspose-slides' library. Please install it: pip install aspose-slides")
        except Exception as e:
            return self.fail_response(f"Failed to export presentation: {str(e)}")
